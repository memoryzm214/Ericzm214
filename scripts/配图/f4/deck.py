# -*- coding: utf-8 -*-
"""把 SVG 版配图重建为 PowerPoint 原生图形的小工具。
所有图元都是可在 PowerPoint 里直接选中、改色、改字的对象，不是位图。"""
from pptx import Presentation
from pptx.util import Pt, Mm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

FONT = 'Microsoft YaHei'
PAGE_W, PAGE_H = 420.0, 297.0        # A3 横向（mm）
MARGIN = 6.0                          # mm


def _rgb(c):
    return RGBColor.from_string(c.lstrip('#').upper())


def _alpha(spPr, tagname, op):
    """给 spPr 下的 solidFill / line 颜色加透明度。"""
    if op >= 1:
        return
    node = spPr.find(qn(tagname))
    if node is None:
        return
    clr = node.find(qn('a:srgbClr'))
    if clr is None:
        return
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(round(op * 100000))))


def _sz(s, fs):
    n = 0.0
    for ch in s:
        n += 1.0 if ord(ch) > 0x2E80 else 0.55
    return n * fs


class Page:
    """一张幻灯片；坐标沿用 SVG 的像素坐标系，内部换算成 EMU。"""

    def __init__(self, prs, W, H):
        self.prs = prs
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        aw = PAGE_W - 2 * MARGIN
        ah = PAGE_H - 2 * MARGIN
        self.k = min(aw / W, ah / H)                 # mm / px
        self.ox = (PAGE_W - W * self.k) / 2
        self.oy = (PAGE_H - H * self.k) / 2
        self.W, self.H = W, H

    # ── 坐标换算 ──
    def X(self, x): return Mm(self.ox + x * self.k)
    def Y(self, y): return Mm(self.oy + y * self.k)
    def L(self, v): return Mm(v * self.k)
    def P(self, v): return Pt(v * self.k * 72 / 25.4)     # px → pt

    # ── 基本图元 ──
    @staticmethod
    def _nostyle(sh):
        st = sh._element.find(qn('p:style'))
        if st is not None:
            sh._element.remove(st)

    def _finish(self, sh, fill, op, stroke, sw, dash=None):
        self._nostyle(sh)
        spPr = sh._element.spPr
        if fill in (None, 'none'):
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = _rgb(fill)
            _alpha(spPr, 'a:solidFill', op)
        if stroke in (None, 'none'):
            sh.line.fill.background()
        else:
            sh.line.color.rgb = _rgb(stroke)
            sh.line.width = self.P(sw)
            if dash:
                sh.line.dash_style = dash
        sh.shadow.inherit = False
        return sh

    def rect(self, x, y, w, h, fill, rx=0, op=1, stroke=None, sw=1, name=None):
        if rx:
            sh = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             self.X(x), self.Y(y), self.L(w), self.L(h))
            try:
                sh.adjustments[0] = min(0.5, rx / max(1e-6, min(w, h)))
            except Exception:
                pass
        else:
            sh = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             self.X(x), self.Y(y), self.L(w), self.L(h))
        if name:
            sh.name = name
        return self._finish(sh, fill, op, stroke, sw)

    def ell(self, cx, cy, rx, ry, fill, op=1, stroke=None, sw=1, name=None):
        sh = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, self.X(cx - rx), self.Y(cy - ry),
                                         self.L(2 * rx), self.L(2 * ry))
        if name:
            sh.name = name
        return self._finish(sh, fill, op, stroke, sw)

    def poly(self, pts, fill, op=1, stroke=None, sw=1, name=None):
        b = self.slide.shapes.build_freeform(self.X(pts[0][0]), self.Y(pts[0][1]))
        b.add_line_segments([(self.X(a), self.Y(c)) for a, c in pts[1:]], close=True)
        sh = b.convert_to_shape()
        if name:
            sh.name = name
        return self._finish(sh, fill, op, stroke, sw)

    def ln(self, x1, y1, x2, y2, color, sw=1, dash=None, arrow=False, cap=None, op=1):
        sh = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                             self.X(x1), self.Y(y1), self.X(x2), self.Y(y2))
        sh.line.color.rgb = _rgb(color)
        sh.line.width = self.P(sw)
        if dash:
            sh.line.dash_style = dash
        lnEl = sh._element.spPr.find(qn('a:ln'))
        if op < 1:
            _alpha(lnEl, 'a:solidFill', op)
        if cap:
            lnEl.set('cap', cap)
        if arrow:
            lnEl.append(parse_xml(
                '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                ' type="triangle" w="med" len="med"/>'))
        sh.shadow.inherit = False
        self._nostyle(sh)
        return sh

    # ── 文字 ──
    def _fmt(self, tf, lines, fs, wt, color, align, op=1):
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        fss = fs if isinstance(fs, (list, tuple)) else [fs] * len(lines)
        wts = wt if isinstance(wt, (list, tuple)) else [wt] * len(lines)
        cols = color if isinstance(color, (list, tuple)) else [color] * len(lines)
        for i, s in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = 1.28
            r = p.add_run()
            r.text = s
            f = r.font
            f.size = self.P(fss[i])
            f.bold = wts[i] >= 600
            f.color.rgb = _rgb(cols[i])
            f.name = FONT
            rPr = r._r.get_or_add_rPr()
            for tag in ('a:latin', 'a:ea', 'a:cs'):
                e = rPr.find(qn(tag))
                if e is None:
                    e = etree.SubElement(rPr, qn(tag))
                e.set('typeface', FONT)
            if op < 1:
                _alpha(rPr, 'a:solidFill', op)

    def txt(self, x, y, s, fs=15, wt=400, color='#31332C', anchor='start', op=1):
        """x,y 为 SVG 的 text 锚点（y 是基线）。"""
        w = _sz(s, fs) * 1.14 + fs * .6
        h = fs * 1.9
        left = {'start': x, 'middle': x - w / 2, 'end': x - w}[anchor]
        align = {'start': PP_ALIGN.LEFT, 'middle': PP_ALIGN.CENTER, 'end': PP_ALIGN.RIGHT}[anchor]
        top = y - fs * .35 - h / 2
        tb = self.slide.shapes.add_textbox(self.X(left), self.Y(top), self.L(w), self.L(h))
        self._fmt(tb.text_frame, [s], fs, wt, color, align, op)
        return tb

    def label(self, sh, lines, fs, wt, color, align=PP_ALIGN.CENTER, op=1):
        """把文字放进图形内部（可随图形一起移动）。"""
        self._fmt(sh.text_frame, lines, fs, wt, color, align, op)
        return sh


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Mm(PAGE_W)
        self.prs.slide_height = Mm(PAGE_H)

    def page(self, W, H):
        return Page(self.prs, W, H)

    def save(self, path):
        self.prs.save(path)
